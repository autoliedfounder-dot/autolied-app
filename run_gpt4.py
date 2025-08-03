import os
import base64
import io
from typing import List, Dict, Optional
import fitz  # PyMuPDF
from PIL import Image
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import re

class SheetMusicToPresentation:
    def __init__(self, openai_api_key: str):
        """
        Initialize with OpenAI API key for GPT-4 Vision.
        """
        self.client = OpenAI(api_key=openai_api_key)
        self.lyrics_data = []
        
    def extract_lyrics_from_pdf(self, pdf_path: str, verbose: bool = True) -> List[Dict[str, any]]:
        """
        Extract lyrics from sheet music PDF using GPT-4 Vision API.
        Returns a list of dictionaries containing lyrics and metadata.
        """
        pdf_document = fitz.open(pdf_path)
        all_page_lyrics = []
        
        # PASS 1: Extract raw lyrics from each page
        for page_num in range(pdf_document.page_count):
            if verbose:
                print(f"Processing page {page_num + 1} of {pdf_document.page_count}...")
            
            # Convert PDF page to image
            page = pdf_document[page_num]
            
            # Start with lower DPI and increase if needed
            dpi = 200  # Start with 200 DPI instead of 300
            max_size = 20 * 1024 * 1024  # 20MB for GPT-4 (larger limit than Claude)
            
            while dpi >= 100:  # Don't go below 100 DPI
                mat = fitz.Matrix(dpi/72, dpi/72)
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")
                
                # Check size
                if len(img_data) <= max_size:
                    break
                
                # If still too large, reduce DPI
                dpi -= 50
                if verbose and dpi >= 100:
                    print(f"  Image too large ({len(img_data)/1024/1024:.1f}MB), reducing to {dpi} DPI...")
            
            # If still too large at 100 DPI, compress using PIL
            if len(img_data) > max_size:
                img_data = self._compress_image(img_data, max_size, verbose)
            
            # Convert to base64 for GPT-4 API
            img_base64 = base64.b64encode(img_data).decode('utf-8')
            
            # Extract lyrics using GPT-4 Vision
            lyrics_on_page = self._extract_lyrics_with_gpt4(img_base64, page_num + 1)
            
            if lyrics_on_page:
                all_page_lyrics.append({
                    'page': page_num + 1,
                    'lyrics': lyrics_on_page
                })
        
        pdf_document.close()
        
        # Combine all pages into a single text for processing
        if verbose:
            print("\nCombining lyrics from all pages...")
        
        combined_lyrics = self._combine_pages_lyrics(all_page_lyrics)
        
        # Post-process and organize lyrics
        self.lyrics_data = self._organize_lyrics(combined_lyrics)
        
        # PASS 2: Reformat lyrics for proper poetic structure
        if verbose:
            print("Reformatting lyrics into proper poetic lines...")
        
        self.lyrics_data = self._reformat_lyrics_poetically(self.lyrics_data)
        
        return self.lyrics_data
    
    def _combine_pages_lyrics(self, all_page_lyrics: List[Dict]) -> List[Dict[str, any]]:
        """
        Combine lyrics from multiple pages into coherent sections.
        This helps avoid breaking phrases across pages.
        """
        combined = []
        current_section = None
        current_text = []
        
        for page_data in all_page_lyrics:
            for lyric in page_data['lyrics']:
                section = lyric.get('section', '')
                text = lyric.get('text', '')
                
                # If section changes, save current and start new
                if section != current_section and current_text:
                    combined.append({
                        'text': '\n'.join(current_text),
                        'section': current_section,
                        'line_number': len(combined) + 1
                    })
                    current_text = []
                
                current_section = section
                if text:
                    current_text.append(text)
        
        # Don't forget the last section
        if current_text:
            combined.append({
                'text': '\n'.join(current_text),
                'section': current_section,
                'line_number': len(combined) + 1
            })
        
        return combined
    
    def _reformat_lyrics_poetically(self, lyrics_data: List[Dict[str, any]]) -> List[Dict[str, any]]:
        """
        Second pass: Use GPT-4 to properly format the extracted lyrics as poetry.
        """
        reformatted_lyrics = []
        
        for lyric in lyrics_data:
            text = lyric.get('text', '')
            section = lyric.get('section', '')
            
            if not text:
                continue
            
            try:
                # First, try to get reference structure from known poems
                reference_structure = self._get_poem_reference_structure(text)
                
                # Ask GPT-4 to reformat the text as poetry
                if reference_structure:
                    prompt = f"""Please format this text as poetry with proper line breaks.

I have found what might be the original poem structure for reference:
{reference_structure}

IMPORTANT: The sheet music may differ from the original poem:
- Lines may be repeated
- Some verses may be omitted
- The order might be changed
- Musical repetitions may be added

Use the reference as a GUIDE for line breaks, but format based on what's ACTUALLY in the extracted text.

Text from sheet music to format:
{text}

Return ONLY the properly formatted text with correct line breaks."""
                else:
                    # Use original prompt if no reference found
                    prompt = f"""Please format this text as poetry with proper line breaks.

CRITICAL RULES:
1. Each poetic line should be a complete thought or phrase
2. Lines that are currently broken mid-phrase should be combined
3. Look for rhyme patterns - rhyming words usually end their lines
4. A typical poetic line is 8-15 words (not too short, not too long)
5. If you see fragments like:
   "my Alice fair," (too short)
   "Her cheek like" (incomplete)
   "the first May rose," (fragment)
   Combine them into: "my Alice fair, Her cheek like the first May rose,"

The text may have incorrect line breaks. Fix them to create proper poetic lines.
Do NOT add or remove any words. Only adjust line breaks.

Text to format:
{text}

Return ONLY the properly formatted poem with correct line breaks."""
                
                response = self.client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {
                            "role": "system",
                            "content": "You are an expert at formatting poetry and song lyrics."
                        },
                        {
                            "role": "user",
                            "content": prompt
                        }
                    ],
                    temperature=0.3,
                    max_tokens=1000
                )
                
                formatted_text = response.choices[0].message.content.strip()
                
                # Remove any markdown formatting (backticks)
                formatted_text = formatted_text.replace('```', '')
                
                # Create new lyric entry with formatted text
                reformatted_lyrics.append({
                    'text': formatted_text,
                    'section': section,
                    'line_number': lyric.get('line_number', 1),
                    'page': lyric.get('page', 1)
                })
                
            except Exception as e:
                print(f"Error reformatting lyrics: {str(e)}")
                # If reformatting fails, keep the original
                reformatted_lyrics.append(lyric)
        
        return reformatted_lyrics
    
    def _get_poem_reference_structure(self, extracted_text: str) -> Optional[str]:
        """
        Try to find the original poem structure for reference.
        Returns the poem with proper line breaks if found, None otherwise.
        """
        try:
            # Get first few lines to identify the poem
            lines = extracted_text.split('\n')
            search_text = ' '.join(lines[:2]) if len(lines) >= 2 else extracted_text[:100]
            
            response = self.client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert on poetry and art songs."
                    },
                    {
                        "role": "user",
                        "content": f"""Can you identify this poem and provide its original structure?

Text excerpt: "{search_text}"

If you recognize this poem, please provide:
1. The title and author
2. The original poem with proper line breaks

If you don't recognize it, just respond with "UNKNOWN".

Note: This is from sheet music, so it might be a fragment or have modifications."""
                    }
                ],
                temperature=0.1,
                max_tokens=1000
            )
            
            reference_text = response.choices[0].message.content.strip()
            
            if "UNKNOWN" in reference_text:
                return None
            else:
                print(f"  Found reference structure for poem")
                return reference_text
                
        except Exception as e:
            print(f"  Could not find reference structure: {str(e)}")
            return None
    
    def _compress_image(self, img_data: bytes, max_size: int, verbose: bool = True) -> bytes:
        """
        Compress image to fit within size limit.
        """
        img = Image.open(io.BytesIO(img_data))
        
        # Convert to RGB if necessary
        if img.mode != 'RGB':
            img = img.convert('RGB')
        
        # Start with 85% quality
        quality = 85
        
        while quality >= 20:
            # Save to bytes with current quality
            output = io.BytesIO()
            img.save(output, format='JPEG', quality=quality, optimize=True)
            compressed_data = output.getvalue()
            
            if len(compressed_data) <= max_size:
                if verbose:
                    print(f"  Compressed to {len(compressed_data)/1024/1024:.1f}MB with {quality}% quality")
                return compressed_data
            
            quality -= 10
        
        # If still too large, resize the image
        if verbose:
            print(f"  Resizing image to fit size limit...")
        
        # Reduce size by 50%
        new_size = (img.width // 2, img.height // 2)
        img = img.resize(new_size, Image.Resampling.LANCZOS)
        
        output = io.BytesIO()
        img.save(output, format='JPEG', quality=70, optimize=True)
        return output.getvalue()
    
    def _extract_lyrics_with_gpt4(self, img_base64: str, page_num: int) -> List[Dict[str, any]]:
        """
        Use GPT-4 Vision API to extract lyrics from a sheet music image.
        """
        try:
            response = self.client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert at reading sheet music and extracting lyrics."
                    },
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": """Analyze this sheet music and extract the lyrics following these rules:

IMPORTANT: When extracting lyrics, keep phrases together! If you see text that continues on the next system, combine it.
For example: "Et sur mon" on one line and "Cœur ouvert, penché," on the next should be extracted as one continuous phrase.

1. STROPHIC SONGS (with numbered verses 1. 2. 3. etc.):
   - Each number continues across ALL systems on the page
   - Combine ALL text from each verse number into complete sentences
   - Track verse 1 through the entire page, then verse 2, etc.

2. MULTILINGUAL SONGS (with translations):
   - ONLY extract the TOP/FIRST language line
   - IGNORE all translation lines below
   
3. STANDARD SONGS:
   - Extract lyrics as they appear
   - PRESERVE POETIC LINE BREAKS - put each poetic line on a new line
   - BUT keep phrases together even if they span multiple systems

Return ONLY a JSON object with this exact structure:
{
    "lyrics": [
        {
            "text": "First line of verse\\nSecond line of verse\\nThird line of verse",
            "section": "verse 1/verse 2/chorus/etc",
            "line_number": 1
        }
    ],
    "song_title": "title if visible",
    "song_type": "strophic/multilingual/standard"
}"""
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{img_base64}",
                                    "detail": "high"  # Use high detail for better accuracy
                                }
                            }
                        ]
                    }
                ],
                max_tokens=2000,
                temperature=0.1  # Lower temperature for more consistent results
            )
            
            # Parse GPT-4's response
            response_text = response.choices[0].message.content
            
            # Extract JSON from response
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if json_match:
                lyrics_data = json.loads(json_match.group())
                
                # Add page number to each lyric entry
                for lyric in lyrics_data.get('lyrics', []):
                    lyric['page'] = page_num
                
                # Print detected type for debugging
                song_type = lyrics_data.get('song_type', 'unknown')
                print(f"  Page {page_num}: Detected as {song_type}")
                
                return lyrics_data.get('lyrics', [])
            else:
                print(f"Could not parse JSON from GPT-4's response on page {page_num}")
                return []
                
        except Exception as e:
            print(f"Error processing page {page_num} with GPT-4: {str(e)}")
            return []
    
    def _organize_lyrics(self, all_lyrics: List[Dict[str, any]]) -> List[Dict[str, any]]:
        """
        Organize and clean up extracted lyrics.
        """
        organized = []
        
        for lyric in all_lyrics:
            text = lyric.get('text', '').strip()
            
            # Skip empty entries
            if text:
                organized.append(lyric)
        
        return organized
    
    def segment_lyrics_for_slides(self, max_lines_per_slide: int = 1) -> List[List[Dict[str, any]]]:
        """
        Segment lyrics into groups for slides.
        For poetry, each line gets its own slide.
        """
        if not self.lyrics_data:
            return []
        
        segments = []
        
        for lyric in self.lyrics_data:
            text = lyric.get('text', '')
            section = lyric.get('section', '')
            
            # Split the text into individual lines
            lines = text.split('\n')
            
            for line in lines:
                line = line.strip()
                if line:  # Skip empty lines
                    # Create a new lyric entry for each line
                    line_entry = {
                        'text': line,
                        'section': section,
                        'line_number': lyric.get('line_number', 1),
                        'page': lyric.get('page', 1)
                    }
                    # Each line gets its own slide
                    segments.append([line_entry])
        
        return segments
    
    def create_powerpoint(self, output_path: str, title: str = "Lyrics Presentation", 
                         max_lines_per_slide: int = 1):
        """
        Create a PowerPoint presentation with the extracted lyrics.
        """
        # Create presentation
        prs = Presentation()
        
        # Set slide size (16:9 ratio)
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        # Add title slide
        self._add_title_slide(prs, title)
        
        # Segment lyrics for slides
        segments = self.segment_lyrics_for_slides(max_lines_per_slide)
        
        # Add lyrics slides
        for i, segment in enumerate(segments):
            lyrics_text = '\n'.join([lyric['text'] for lyric in segment])
            section = segment[0].get('section', '') if segment else ''
            
            self._add_lyrics_slide(prs, lyrics_text, section, i + 1, len(segments))
        
        # Save presentation
        prs.save(output_path)
        
        return len(segments)  # Return number of slides created
    
    def _add_title_slide(self, prs: Presentation, title: str):
        """
        Add a title slide to the presentation.
        """
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background
        self._add_background(slide, RGBColor(0, 0, 0))
        
        # Add title text box
        left = Inches(2)
        top = Inches(3.5)
        width = Inches(12)
        height = Inches(2)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.add_paragraph()
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        
        # Style the title
        font = p.font
        font.size = Pt(60)
        font.name = 'Arial'
        font.bold = True
        font.color.rgb = RGBColor(255, 255, 255)
    
    def _add_lyrics_slide(self, prs: Presentation, lyrics: str, section: str, 
                         slide_num: int, total_slides: int):
        """
        Add a lyrics slide with centered text and optional section indicator.
        """
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background
        self._add_background(slide, RGBColor(0, 0, 0))
        
        # Add section indicator if present (smaller text at top)
        if section:
            section_left = Inches(1)
            section_top = Inches(0.5)
            section_width = Inches(14)
            section_height = Inches(1)
            
            section_box = slide.shapes.add_textbox(section_left, section_top, 
                                                   section_width, section_height)
            section_frame = section_box.text_frame
            
            p = section_frame.add_paragraph()
            p.text = f"[{section.upper()}]"
            p.alignment = PP_ALIGN.CENTER
            
            font = p.font
            font.size = Pt(24)
            font.name = 'Arial'
            font.italic = True
            font.color.rgb = RGBColor(200, 200, 200)  # Light gray
        
        # Add main lyrics text box
        left = Inches(1)
        top = Inches(2.5) if section else Inches(2)
        width = Inches(14)
        height = Inches(5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.add_paragraph()
        p.text = lyrics
        p.alignment = PP_ALIGN.CENTER
        
        # Style the lyrics
        font = p.font
        font.size = Pt(48)
        font.name = 'Arial'
        font.bold = True
        font.color.rgb = RGBColor(255, 255, 255)
        
        # Add slide number (bottom right)
        self._add_slide_number(slide, slide_num, total_slides)
    
    def _add_slide_number(self, slide, current: int, total: int):
        """
        Add slide number indicator.
        """
        left = Inches(14.5)
        top = Inches(8.5)
        width = Inches(1.5)
        height = Inches(0.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        p = text_frame.add_paragraph()
        p.text = f"{current}/{total}"
        p.alignment = PP_ALIGN.RIGHT
        
        font = p.font
        font.size = Pt(14)
        font.name = 'Arial'
        font.color.rgb = RGBColor(150, 150, 150)
    
    def _add_background(self, slide, color: RGBColor):
        """
        Add a colored background to a slide.
        """
        left = top = 0
        # Get presentation from slide
        prs = slide.part.package.presentation_part.presentation
        width = prs.slide_width
        height = prs.slide_height
        
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = color
        
        # Send to back
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.insert(2, shape._element)
    
    def export_lyrics_to_text(self, output_path: str):
        """
        Export extracted lyrics to a text file for review/editing.
        """
        with open(output_path, 'w', encoding='utf-8') as f:
            current_section = None
            
            for lyric in self.lyrics_data:
                section = lyric.get('section', '')
                
                # Add section header if it changes
                if section and section != current_section:
                    f.write(f"\n[{section.upper()}]\n")
                    current_section = section
                
                f.write(f"{lyric['text']}\n")

# Main processing function
def process_sheet_music_with_gpt4(pdf_path: str, output_pptx: str, api_key: str, 
                                  title: Optional[str] = None,
                                  max_lines_per_slide: int = 1,
                                  export_text: bool = True):
    """
    Complete pipeline to process sheet music PDF and create PowerPoint using GPT-4 Vision.
    
    Args:
        pdf_path: Path to the sheet music PDF
        output_pptx: Path for the output PowerPoint file
        api_key: OpenAI API key
        title: Title for the presentation (optional)
        max_lines_per_slide: Maximum lines of lyrics per slide
        export_text: Whether to also export lyrics as text file
    """
    print("Initializing sheet music processor with GPT-4 Vision...")
    processor = SheetMusicToPresentation(api_key)
    
    # Extract lyrics
    print(f"Extracting lyrics from {pdf_path}...")
    lyrics_data = processor.extract_lyrics_from_pdf(pdf_path)
    
    if not lyrics_data:
        print("No lyrics found in the PDF.")
        return
    
    print(f"Successfully extracted {len(lyrics_data)} lyric lines")
    
    # Export to text file if requested
    if export_text:
        text_path = output_pptx.replace('.pptx', '_lyrics.txt')
        processor.export_lyrics_to_text(text_path)
        print(f"Lyrics exported to: {text_path}")
    
    # Determine title
    if not title:
        # Try to find title from extracted data
        for lyric in lyrics_data:
            if lyric.get('song_title'):
                title = lyric['song_title']
                break
        if not title:
            title = "Lyrics Presentation"
    
    # Create PowerPoint
    print(f"Creating PowerPoint presentation...")
    num_slides = processor.create_powerpoint(output_pptx, title, max_lines_per_slide)
    
    print(f"✅ PowerPoint created successfully!")
    print(f"   - Output file: {output_pptx}")
    print(f"   - Total slides: {num_slides + 1} (including title slide)")
    
    return processor

# Example usage
if __name__ == "__main__":
    # Configuration
    PDF_PATH = "your_sheet_music.pdf"
    OUTPUT_PATH = "lyrics_presentation.pptx"
    API_KEY = "your-openai-api-key"  # Replace with your OpenAI API key
    
    # Process the sheet music
    processor = process_sheet_music_with_gpt4(
        pdf_path=PDF_PATH,
        output_pptx=OUTPUT_PATH,
        api_key=API_KEY,
        title="My Song Lyrics",  # Optional: will try to detect from PDF
        max_lines_per_slide=2,   # Adjust based on your preference
        export_text=True         # Also creates a text file with lyrics
    )